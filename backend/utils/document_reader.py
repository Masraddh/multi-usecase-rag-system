"""
Document Reader Utility Module for RAG Pipeline.

Provides robust document reading, cleaning, validation, and text extraction
for PDF, DOCX, TXT, and Markdown files.

Primary PDF Loader: PyMuPDF (fitz)
Fallback PDF Loader 1: pdfplumber
Fallback PDF Loader 2: pypdf
"""

import os
import io
import re
import csv
from typing import Tuple, Optional, Union, List, BinaryIO

# Primary PDF Loader
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

# Fallback 1 PDF Loader
try:
    import pdfplumber
except ImportError:
    pdfplumber = None

# Fallback 2 PDF Loader
try:
    import pypdf
except ImportError:
    pypdf = None

# DOCX Loader
try:
    import docx
except ImportError:
    docx = None

# PPTX Loader
try:
    import pptx
except ImportError:
    pptx = None


def clean_text(text: str) -> str:
    """
    Cleans raw text extracted from documents:
    - Removes null bytes and non-printable control characters
    - Normalizes Windows/Mac line endings to Unix \n
    - Replaces multiple consecutive blank lines with double newlines
    - Trims leading/trailing space within lines while preserving paragraph breaks
    """
    if not text:
        return ""
    
    text = text.replace("\x00", "")
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"\n{3,}", "\n\n", text)
    lines = [re.sub(r"[ \t]+", " ", line).strip() for line in text.split("\n")]
    cleaned = "\n".join(lines).strip()
    return cleaned


def validate_document(text: str) -> bool:
    """
    Validates that extracted text is non-empty and contains readable words.
    Returns True if valid, False otherwise.
    """
    if not text or not isinstance(text, str):
        return False
    words = text.strip().split()
    return len(words) > 0


def _get_bytes(file_input: Union[str, bytes, BinaryIO]) -> bytes:
    """
    Helper function to convert file paths, byte arrays, or file-like objects into raw bytes.
    """
    if isinstance(file_input, bytes):
        return file_input
    elif isinstance(file_input, str):
        if os.path.exists(file_input):
            with open(file_input, "rb") as f:
                return f.read()
        else:
            raise ValueError(f"File path not found: '{file_input}'")
    elif hasattr(file_input, "read"):
        content = file_input.read()
        if isinstance(content, str):
            return content.encode("utf-8")
        return content
    else:
        raise ValueError("Invalid file input type. Expected file path (str), bytes, or file stream.")


import traceback


def read_pdf(file_input: Union[str, bytes, BinaryIO], filename: Optional[str] = None) -> Tuple[str, int]:
    """
    Reads all pages from a PDF file using PyMuPDF (fitz), with fallback to pdfplumber and pypdf.
    Logs step-by-step page metrics and unmasked Python tracebacks on failure.
    Returns: (cleaned_combined_text, page_count)
    """
    file_bytes = _get_bytes(file_input)
    file_size_kb = round(len(file_bytes) / 1024, 2)

    print("=" * 80, flush=True)
    print(f"[PDF READER ENTRY DEBUG LOG]", flush=True)
    print(f"- Received File: {filename or 'uploaded_file.pdf'}", flush=True)
    print(f"- File Size: {file_size_kb} KB ({len(file_bytes)} bytes)", flush=True)
    print(f"- Bytes Stream Valid: {'YES' if len(file_bytes) > 0 else 'NO'}", flush=True)
    print("=" * 80, flush=True)

    if not file_bytes:
        raise ValueError("Empty PDF file payload received (0 bytes). Please select a valid PDF file.")

    pages_text = []
    page_count = 0
    errors = []

    # 1. Primary Loader: PyMuPDF (fitz)
    if fitz:
        try:
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            if doc.is_encrypted:
                raise ValueError("This PDF is password-protected or encrypted. Please upload an unencrypted PDF file.")
            
            page_count = len(doc)
            print(f"[PyMuPDF] Opened PDF successfully. Total Pages: {page_count}", flush=True)

            for i, page in enumerate(doc, 1):
                text = page.get_text("text") or page.get_text()
                txt_clean = text.strip() if text else ""
                chars = len(txt_clean)
                words = len(txt_clean.split())
                snippet = txt_clean[:100].replace("\n", " ") if txt_clean else "<EMPTY_PAGE>"
                print(f"  -> Page {i}: Characters={chars}, Words={words}, Preview=\"{snippet}...\"", flush=True)
                if txt_clean:
                    pages_text.append(txt_clean)
            doc.close()
        except ValueError as ve:
            raise ve
        except Exception as e:
            tb = traceback.format_exc()
            print(f"[PyMuPDF ERROR] Failed to parse PDF pages:\n{tb}", flush=True)
            errors.append(f"PyMuPDF error: {str(e)}")
            pages_text = []

    # 2. Fallback 1: pdfplumber
    if not pages_text and pdfplumber:
        try:
            print("[PDF FALLBACK] Attempting pdfplumber extraction...", flush=True)
            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                page_count = len(pdf.pages)
                for i, page in enumerate(pdf.pages, 1):
                    txt = page.extract_text()
                    txt_clean = txt.strip() if txt else ""
                    chars = len(txt_clean)
                    words = len(txt_clean.split())
                    snippet = txt_clean[:100].replace("\n", " ") if txt_clean else "<EMPTY_PAGE>"
                    print(f"  -> [pdfplumber] Page {i}: Characters={chars}, Words={words}, Preview=\"{snippet}...\"", flush=True)
                    if txt_clean:
                        pages_text.append(txt_clean)
        except ValueError as ve:
            raise ve
        except Exception as e:
            tb = traceback.format_exc()
            print(f"[pdfplumber ERROR] Failed to parse PDF:\n{tb}", flush=True)
            errors.append(f"pdfplumber error: {str(e)}")
            pages_text = []

    # 3. Fallback 2: pypdf
    if not pages_text and pypdf:
        try:
            print("[PDF FALLBACK] Attempting pypdf extraction...", flush=True)
            reader = pypdf.PdfReader(io.BytesIO(file_bytes))
            if reader.is_encrypted:
                raise ValueError("This PDF is password-protected or encrypted. Please upload an unencrypted PDF file.")
            
            page_count = len(reader.pages)
            for i, page in enumerate(reader.pages, 1):
                txt = page.extract_text()
                txt_clean = txt.strip() if txt else ""
                chars = len(txt_clean)
                words = len(txt_clean.split())
                snippet = txt_clean[:100].replace("\n", " ") if txt_clean else "<EMPTY_PAGE>"
                print(f"  -> [pypdf] Page {i}: Characters={chars}, Words={words}, Preview=\"{snippet}...\"", flush=True)
                if txt_clean:
                    pages_text.append(txt_clean)
        except ValueError as ve:
            raise ve
        except Exception as e:
            tb = traceback.format_exc()
            print(f"[pypdf ERROR] Failed to parse PDF:\n{tb}", flush=True)
            errors.append(f"pypdf error: {str(e)}")
            pages_text = []

    if not pages_text or page_count == 0:
        if fitz:
            try:
                doc = fitz.open(stream=file_bytes, filetype="pdf")
                if doc.is_encrypted:
                    doc.close()
                    raise ValueError("This PDF is password-protected or encrypted. Please upload an unencrypted PDF file.")
                has_imgs = any(len(page.get_images()) > 0 for page in doc)
                doc.close()
                if has_imgs:
                    raise ValueError("This PDF contains scanned images without selectable text. Please upload a PDF with embedded text or a DOCX/TXT file.")
            except ValueError as ve:
                raise ve
            except Exception:
                pass
        err_msg = "; ".join(errors) if errors else "No readable text found across PDF pages."
        raise ValueError(f"Unable to read this PDF document: {err_msg}")

    combined_text = "\n\n".join(pages_text)
    cleaned = clean_text(combined_text)

    if not validate_document(cleaned):
        raise ValueError("Extracted PDF text is empty or invalid. Ensure the PDF is not a scanned image.")

    return cleaned, max(page_count, 1)


def read_docx(file_input: Union[str, bytes, BinaryIO]) -> Tuple[str, int]:
    """
    Reads text and estimated page count from a Word (.docx) document.
    Extracts text from paragraphs, tables, text boxes (w:txbxContent), headers, footers, and XML nodes.
    """
    if not docx:
        raise ValueError("python-docx library is required to read .docx files.")

    file_bytes = _get_bytes(file_input)
    try:
        doc = docx.Document(io.BytesIO(file_bytes))
        paragraphs = []

        # 1. Main Document Paragraphs
        for p in doc.paragraphs:
            if p.text and p.text.strip():
                paragraphs.append(p.text.strip())

        # 2. Main Document Tables (including cell paragraphs)
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    paragraphs.append(row_text)

        # 3. Text Boxes & Drawing Shapes (w:txbxContent)
        try:
            for txbx in doc.element.xpath("//w:txbxContent"):
                texts = [node.text.strip() for node in txbx.xpath(".//w:t") if node.text and node.text.strip()]
                if texts:
                    paragraphs.append(" ".join(texts))
        except Exception:
            pass

        # 4. Headers & Footers
        try:
            for section in doc.sections:
                if section.header:
                    for hp in section.header.paragraphs:
                        if hp.text and hp.text.strip():
                            paragraphs.append(hp.text.strip())
                if section.footer:
                    for fp in section.footer.paragraphs:
                        if fp.text and fp.text.strip():
                            paragraphs.append(fp.text.strip())
        except Exception:
            pass

        # 5. XML Fallback: if paragraphs list is still empty, extract all <w:t> text nodes
        if not paragraphs:
            try:
                xml_texts = doc.element.xpath("//w:t")
                raw_xml_str = " ".join([t.text.strip() for t in xml_texts if t.text and t.text.strip()])
                if raw_xml_str:
                    paragraphs.append(raw_xml_str)
            except Exception:
                pass

        full_text = "\n\n".join(paragraphs)
        cleaned = clean_text(full_text)
        words = len(cleaned.split())
        page_count = max(1, (words + 349) // 350)

        if not validate_document(cleaned):
            raise ValueError("Unable to read this document. Please upload a valid PDF, DOCX, TXT or Markdown file.")

        return cleaned, page_count
    except ValueError as ve:
        raise ve
    except Exception as e:
        raise ValueError(f"Unable to read this document. Please upload a valid PDF, DOCX, TXT or Markdown file.")


def read_txt(file_input: Union[str, bytes, BinaryIO]) -> Tuple[str, int]:
    """
    Reads plain text (.txt) file contents.
    """
    file_bytes = _get_bytes(file_input)
    try:
        text = file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        try:
            text = file_bytes.decode("latin-1")
        except Exception:
            raise ValueError("Unable to decode text file. Ensure it is UTF-8 encoded.")

    cleaned = clean_text(text)
    words = len(cleaned.split())
    page_count = max(1, (words + 349) // 350)

    if not validate_document(cleaned):
        raise ValueError("Unable to read this document. Please upload a valid PDF, DOCX, TXT or Markdown file.")

    return cleaned, page_count


def read_markdown(file_input: Union[str, bytes, BinaryIO]) -> Tuple[str, int]:
    """
    Reads Markdown (.md) file contents.
    """
    return read_txt(file_input)


def extract_text(
    file_input: Union[str, bytes, BinaryIO],
    filename: Optional[str] = None
) -> Tuple[str, int, int, int]:
    """
    Central router function to read documents, clean text, validate content, and return statistics.
    Returns: (cleaned_text, page_count, word_count, char_count)
    """
    file_bytes = _get_bytes(file_input)
    if not file_bytes:
        raise ValueError("Unable to read this document. Please upload a valid PDF, DOCX, TXT or Markdown file.")

    # Infer extension
    ext = ""
    if filename:
        ext = os.path.splitext(filename)[1].lower()
    elif isinstance(file_input, str):
        ext = os.path.splitext(file_input)[1].lower()

    if not ext:
        ext = ".pdf"  # Default assumption for binary upload stream if un-specified

    if ext == ".pdf":
        cleaned_text, page_count = read_pdf(file_bytes)
    elif ext in [".docx", ".doc"]:
        cleaned_text, page_count = read_docx(file_bytes)
    elif ext == ".txt":
        cleaned_text, page_count = read_txt(file_bytes)
    elif ext == ".md":
        cleaned_text, page_count = read_markdown(file_bytes)
    elif ext == ".pptx":
        if not pptx:
            raise ValueError("python-pptx library is required to read PowerPoint files.")
        prs = pptx.Presentation(io.BytesIO(file_bytes))
        slide_texts = []
        for i, slide in enumerate(prs.slides, 1):
            runs = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if runs:
                slide_texts.append(f"--- Slide {i} ---\n" + "\n".join(runs))
        cleaned_text = clean_text("\n\n".join(slide_texts))
        page_count = len(prs.slides)
    elif ext == ".csv":
        content_str = file_bytes.decode("utf-8", errors="ignore")
        reader = csv.reader(io.StringIO(content_str))
        rows = list(reader)
        header = rows[0] if rows else []
        formatted = []
        for i, r in enumerate(rows[1:], 1):
            items = [f"{header[c] if c < len(header) else f'Col_{c+1}'}: {val.strip()}" for c, val in enumerate(r) if val.strip()]
            if items:
                formatted.append(f"Record {i}: " + "; ".join(items))
        cleaned_text = clean_text("\n".join(formatted))
        page_count = max(1, (len(cleaned_text.split()) + 349) // 350)
    elif ext in [".py", ".json", ".js", ".ts", ".html", ".css", ".yaml", ".yml"]:
        cleaned_text, page_count = read_txt(file_bytes)
    else:
        raise ValueError(
            f"Unsupported file format '{ext}'. Supported formats: PDF (.pdf), Word (.docx), Text (.txt), Markdown (.md)."
        )

    if not validate_document(cleaned_text):
        raise ValueError("Unable to read this document. Please upload a valid PDF, DOCX, TXT or Markdown file.")

    word_count = len(cleaned_text.split())
    char_count = len(cleaned_text)

    preview_500 = cleaned_text[:500].replace("\n", " ")
    print("=" * 80, flush=True)
    print(f"[DOCUMENT EXTRACTION LOG]", flush=True)
    print(f"- Document Name: {filename or 'custom_document'}", flush=True)
    print(f"- Pages Read: {page_count}", flush=True)
    print(f"- Characters Extracted: {char_count}", flush=True)
    print(f"- Words Extracted: {word_count}", flush=True)
    print(f"- First 500 Characters Preview: \"{preview_500}...\"", flush=True)
    print("=" * 80, flush=True)

    return cleaned_text, page_count, word_count, char_count


class DocumentReader:
    """
    Class wrapper mapping static methods to document reading utilities.
    """
    read_pdf = staticmethod(read_pdf)
    read_docx = staticmethod(read_docx)
    read_txt = staticmethod(read_txt)
    read_markdown = staticmethod(read_markdown)
    extract_text = staticmethod(extract_text)
    clean_text = staticmethod(clean_text)
    validate_document = staticmethod(validate_document)
