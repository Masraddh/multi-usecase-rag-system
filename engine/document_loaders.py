import os
import io
import re
import csv
from typing import Tuple, Optional

# Primary PDF Loader
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

# Fallback 1 PDF Loader
try:
    import pdfplumber
except ImportError:
    pdfplumber = None

# Fallback 2 PDF Loader
try:
    import pypdf
except ImportError:
    pypdf = None

# DOCX Loader
try:
    import docx
except ImportError:
    docx = None

# PPTX Loader
try:
    import pptx
except ImportError:
    pptx = None


def clean_text(text: str) -> str:
    """
    Cleans raw text extracted from documents:
    - Removes null bytes and non-printable control characters
    - Normalizes Windows/Mac line endings to Unix \n
    - Replaces multiple consecutive blank lines with double newlines
    - Trims leading/trailing space within lines while preserving paragraph breaks
    """
    if not text:
        return ""
    
    text = text.replace("\x00", "")
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"\n{3,}", "\n\n", text)
    lines = [re.sub(r"[ \t]+", " ", line).strip() for line in text.split("\n")]
    cleaned = "\n".join(lines).strip()
    return cleaned


def load_pdf(file_bytes: bytes) -> Tuple[str, int]:
    """
    Extracts text and page count from a PDF file.
    Primary: PyMuPDF (fitz)
    Fallback 1: pdfplumber
    Fallback 2: pypdf
    """
    if not file_bytes:
        raise ValueError("Unable to read this PDF. Please upload a valid, non-password-protected PDF.")

    pages_text = []
    page_count = 0

    # 1. Primary Loader: PyMuPDF (fitz)
    if fitz:
        try:
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            if doc.is_encrypted:
                raise ValueError("Unable to read this PDF. Please upload a valid, non-password-protected PDF.")
            
            page_count = len(doc)
            for page in doc:
                text = page.get_text("text") or page.get_text()
                if text and text.strip():
                    pages_text.append(text.strip())
            doc.close()
        except ValueError as ve:
            raise ve
        except Exception:
            pages_text = []

    # 2. Fallback 1: pdfplumber
    if not pages_text and pdfplumber:
        try:
            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                page_count = len(pdf.pages)
                for page in pdf.pages:
                    txt = page.extract_text()
                    if txt and txt.strip():
                        pages_text.append(txt.strip())
        except ValueError as ve:
            raise ve
        except Exception:
            pages_text = []

    # 3. Fallback 2: pypdf
    if not pages_text and pypdf:
        try:
            reader = pypdf.PdfReader(io.BytesIO(file_bytes))
            if reader.is_encrypted:
                raise ValueError("Unable to read this PDF. Please upload a valid, non-password-protected PDF.")
            
            page_count = len(reader.pages)
            for page in reader.pages:
                txt = page.extract_text()
                if txt and txt.strip():
                    pages_text.append(txt.strip())
        except ValueError as ve:
            raise ve
        except Exception:
            pages_text = []

    if not pages_text or page_count == 0:
        raise ValueError("Unable to read this PDF. Please upload a valid, non-password-protected PDF.")

    combined_text = "\n\n".join(pages_text)
    cleaned = clean_text(combined_text)

    if not cleaned or len(cleaned.split()) == 0:
        raise ValueError("Unable to read this PDF. Please upload a valid, non-password-protected PDF.")

    return cleaned, max(page_count, 1)


def load_docx(file_bytes: bytes) -> Tuple[str, int]:
    """
    Extracts text and estimated page count from a Word (.docx) document.
    """
    if not docx:
        raise ValueError("python-docx library is required to read .docx files.")

    try:
        doc = docx.Document(io.BytesIO(file_bytes))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    paragraphs.append(row_text)

        full_text = "\n\n".join(paragraphs)
        cleaned = clean_text(full_text)
        words = len(cleaned.split())
        page_count = max(1, (words + 349) // 350)

        return cleaned, page_count
    except Exception as e:
        raise ValueError(f"Unable to read DOCX document: {str(e)}")


def load_txt(file_bytes: bytes) -> Tuple[str, int]:
    """
    Reads plain text (.txt) file contents.
    """
    try:
        text = file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        try:
            text = file_bytes.decode("latin-1")
        except Exception:
            raise ValueError("Unable to decode text file. Ensure it is UTF-8 encoded.")

    cleaned = clean_text(text)
    words = len(cleaned.split())
    page_count = max(1, (words + 349) // 350)
    return cleaned, page_count


def load_markdown(file_bytes: bytes) -> Tuple[str, int]:
    """
    Reads Markdown (.md) file contents.
    """
    return load_txt(file_bytes)


def load_pptx(file_bytes: bytes) -> Tuple[str, int]:
    """
    Extracts text and slide count from a PowerPoint (.pptx) presentation.
    """
    if not pptx:
        raise ValueError("python-pptx library is not installed. Please install python-pptx to process PowerPoint files.")

    try:
        prs = pptx.Presentation(io.BytesIO(file_bytes))
        slide_texts = []
        slide_count = len(prs.slides)

        for i, slide in enumerate(prs.slides, 1):
            text_runs = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_runs.append(shape.text.strip())
            if text_runs:
                slide_texts.append(f"--- Slide {i} ---\n" + "\n".join(text_runs))

        full_text = "\n\n".join(slide_texts)
        cleaned = clean_text(full_text)
        return cleaned, max(slide_count, 1)
    except Exception as e:
        raise ValueError(f"Unable to parse PowerPoint file: {str(e)}")


def load_csv(file_bytes: bytes) -> Tuple[str, int]:
    """
    Reads CSV file content and formats rows into readable key-value sentences.
    """
    try:
        content_str = file_bytes.decode("utf-8", errors="ignore")
        reader = csv.reader(io.StringIO(content_str))
        rows = list(reader)
        if not rows:
            return "", 1

        header = rows[0]
        formatted_rows = []
        for i, row in enumerate(rows[1:], 1):
            row_items = []
            for col_idx, val in enumerate(row):
                col_name = header[col_idx] if col_idx < len(header) else f"Col_{col_idx+1}"
                if val.strip():
                    row_items.append(f"{col_name}: {val.strip()}")
            if row_items:
                formatted_rows.append(f"Record {i}: " + "; ".join(row_items))

        full_text = "\n".join(formatted_rows)
        cleaned = clean_text(full_text)
        words = len(cleaned.split())
        page_count = max(1, (words + 349) // 350)
        return cleaned, page_count
    except Exception as e:
        raise ValueError(f"Unable to read CSV file: {str(e)}")


def extract_text(filename: str, file_bytes: bytes) -> Tuple[str, int, int, int]:
    """
    Central router function to validate, load, clean text, and return stats.
    Returns: (cleaned_text, page_count, word_count, char_count)
    """
    if not file_bytes:
        raise ValueError("Uploaded file is empty (0 bytes). Please upload a valid document.")

    ext = os.path.splitext(filename)[1].lower()

    if ext == ".pdf":
        cleaned_text, page_count = load_pdf(file_bytes)
    elif ext in [".docx", ".doc"]:
        cleaned_text, page_count = load_docx(file_bytes)
    elif ext == ".txt":
        cleaned_text, page_count = load_txt(file_bytes)
    elif ext == ".md":
        cleaned_text, page_count = load_markdown(file_bytes)
    elif ext == ".pptx":
        cleaned_text, page_count = load_pptx(file_bytes)
    elif ext == ".csv":
        cleaned_text, page_count = load_csv(file_bytes)
    elif ext in [".py", ".json", ".js", ".ts", ".html", ".css", ".yaml", ".yml"]:
        cleaned_text, page_count = load_txt(file_bytes)
    else:
        raise ValueError(
            f"Unsupported file format '{ext}'. Supported formats: PDF (.pdf), Word (.docx), Text (.txt), Markdown (.md), PowerPoint (.pptx), CSV (.csv)."
        )

    word_count = len(cleaned_text.split())
    char_count = len(cleaned_text)

    if word_count == 0 or char_count == 0:
        raise ValueError("Unable to read this PDF. Please upload a valid, non-password-protected PDF.")

    return cleaned_text, page_count, word_count, char_count
